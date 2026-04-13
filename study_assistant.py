import os
import json
import re
import threading
import tkinter as tk
from tkinter import messagebox, filedialog
import customtkinter as ctk
from PyPDF2 import PdfReader
from pptx import Presentation
import google.generativeai as genai

# --- UI Theme Settings ---
ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("blue")

class StudyAssistantApp(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("🎓 AI 스마트 강의 요약 및 예상 문제 생성기")
        self.geometry("1100x850")

        # 데이터 변수
        self.api_key = ""
        self.extracted_text = ""
        self.quiz_data = []
        self.user_answers = []
        self.explanation_labels = []

        self.setup_ui()

    def setup_ui(self):
        # 그리드 설정
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)

        # 1. 사이드바 (설정 영역)
        self.sidebar = ctk.CTkFrame(self, width=250, corner_radius=0)
        self.sidebar.grid(row=0, column=0, sticky="nsew")
        self.sidebar.grid_rowconfigure(6, weight=1)

        ctk.CTkLabel(self.sidebar, text="⚙️ 설정 및 파일", font=ctk.CTkFont(size=20, weight="bold")).grid(row=0, column=0, padx=20, pady=20)

        # API Key 입력
        ctk.CTkLabel(self.sidebar, text="Gemini API Key:", font=ctk.CTkFont(size=12)).grid(row=1, column=0, padx=20, pady=(10, 0), sticky="w")
        self.api_entry = ctk.CTkEntry(self.sidebar, placeholder_text="API 키를 입력하세요", show="*")
        self.api_entry.grid(row=2, column=0, padx=20, pady=(5, 20), sticky="ew")

        # 파일 선택 버튼들
        self.btn_load_pdf = ctk.CTkButton(self.sidebar, text="📄 PDF 파일 불러오기", command=self.load_pdf, fg_color="#3498db", hover_color="#2980b9")
        self.btn_load_pdf.grid(row=3, column=0, padx=20, pady=10, sticky="ew")

        self.btn_load_pptx = ctk.CTkButton(self.sidebar, text="📊 PPTX 파일 불러오기", command=self.load_pptx, fg_color="#e67e22", hover_color="#d35400")
        self.btn_load_pptx.grid(row=4, column=0, padx=20, pady=10, sticky="ew")

        self.appearance_mode_label = ctk.CTkLabel(self.sidebar, text="테마 설정:", anchor="w")
        self.appearance_mode_label.grid(row=7, column=0, padx=20, pady=(10, 0))
        self.appearance_mode_optionemenu = ctk.CTkOptionMenu(self.sidebar, values=["Light", "Dark", "System"], command=self.change_appearance_mode_event)
        self.appearance_mode_optionemenu.set("Dark")
        self.appearance_mode_optionemenu.grid(row=8, column=0, padx=20, pady=(10, 20))

        # 2. 메인 화면 영역 (탭 뷰)
        self.tabview = ctk.CTkTabview(self, corner_radius=10)
        self.tabview.grid(row=0, column=1, padx=20, pady=20, sticky="nsew")
        self.tabview.add("입력 내용")
        self.tabview.add("5줄 요약")
        self.tabview.add("예상 문제")

        # [탭 1: 입력 내용]
        self.input_textbox = ctk.CTkTextbox(self.tabview.tab("입력 내용"), font=ctk.CTkFont(size=13))
        self.input_textbox.pack(padx=10, pady=10, fill="both", expand=True)
        
        self.process_btn = ctk.CTkButton(self.tabview.tab("입력 내용"), text="✨ AI 요약 및 문제 생성 시작", 
                                        height=45, font=ctk.CTkFont(size=15, weight="bold"),
                                        command=self.start_ai_process)
        self.process_btn.pack(padx=10, pady=10, fill="x")

        # [탭 2: 5줄 요약]
        self.summary_textbox = ctk.CTkTextbox(self.tabview.tab("5줄 요약"), font=ctk.CTkFont(size=14), state="disabled")
        self.summary_textbox.pack(padx=20, pady=20, fill="both", expand=True)

        # [탭 3: 예상 문제]
        self.quiz_scroll = ctk.CTkScrollableFrame(self.tabview.tab("예상 문제"), label_text="강의 내용 기반 예상 문제")
        self.quiz_scroll.pack(padx=10, pady=10, fill="both", expand=True)
        
        self.grade_btn = ctk.CTkButton(self.tabview.tab("예상 문제"), text="✅ 정답 확인 및 채점", 
                                      state="disabled", command=self.grade_quiz)
        self.grade_btn.pack(padx=10, pady=10, fill="x")

    def change_appearance_mode_event(self, new_appearance_mode: str):
        ctk.set_appearance_mode(new_appearance_mode)

    def load_pdf(self):
        file_path = filedialog.askopenfilename(filetypes=[("PDF Files", "*.pdf")])
        if file_path:
            try:
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                self.input_textbox.delete("1.0", tk.END)
                self.input_textbox.insert("1.0", text)
                messagebox.showinfo("성공", "PDF 텍스트 추출 완료!")
            except Exception as e:
                messagebox.showerror("오류", f"PDF를 읽을 수 없습니다: {e}")

    def load_pptx(self):
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            try:
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                self.input_textbox.delete("1.0", tk.END)
                self.input_textbox.insert("1.0", text)
                messagebox.showinfo("성공", "PPTX 텍스트 추출 완료!")
            except Exception as e:
                messagebox.showerror("오류", f"PPTX를 읽을 수 없습니다: {e}")

    def start_ai_process(self):
        content = self.input_textbox.get("1.0", tk.END).strip()
        api_key = self.api_entry.get().strip()

        if len(content) < 100:
            messagebox.showwarning("입력 부족", "분석할 내용이 너무 적습니다. 최소 100자 이상 입력해주세요.")
            return
        
        if not api_key:
            messagebox.showwarning("API 키 누락", "Gemini API 키를 입력해주세요. (없으면 샘플 모드로 동작합니다)")
            # API 키가 없으면 데모 모드로 전환 여부 확인
            if not messagebox.askyesno("데모 모드", "API 키 없이 샘플 데이터로 시연하시겠습니까?"):
                return

        self.process_btn.configure(state="disabled", text="⏳ AI 분석 중...")
        threading.Thread(target=self.run_gemini, args=(content, api_key), daemon=True).start()

    def run_gemini(self, content, api_key):
        if not api_key:
            # 샘플 데이터 사용
            import time
            time.sleep(1.5)
            sample_data = {
                "summary": [
                    "1. 인공지능의 기본 원리는 데이터를 기반으로 패턴을 학습하는 것입니다.",
                    "2. 머신러닝은 지도 학습, 비지도 학습, 강화 학습 등 세 가지 주요 유형으로 나뉩니다.",
                    "3. 딥러닝은 인간의 뇌 구조를 모방한 인공 신경망을 사용하여 복잡한 문제를 해결합니다.",
                    "4. 대규모 언어 모델(LLM)은 현대 AI의 핵심 기술로 자리 잡고 있습니다.",
                    "5. AI의 발전은 의료, 자동화, 창작 등 다양한 산업 분야에 혁신을 가져오고 있습니다."
                ],
                "quizzes": [
                    {
                        "question": "다음 중 학습 데이터에 정답(Label)이 포함된 학습 방식은 무엇인가요?",
                        "options": ["비지도 학습", "지도 학습", "강화 학습", "자율 학습"],
                        "answer": 2,
                        "explanation": "지도 학습(Supervised Learning)은 입력 데이터와 함께 정답 라벨을 제공하여 모델을 학습시키는 방식입니다."
                    },
                    {
                        "question": "인간의 뇌 구조를 모방한 인공지능 기술의 명칭으로 옳은 것은?",
                        "options": ["선형 회귀", "결정 트리", "인공 신경망(Neural Networks)", "K-평균 알고리즘"],
                        "answer": 3,
                        "explanation": "인공 신경망은 생물학적인 뇌의 뉴런 연결 구조를 본떠 만든 알고리즘입니다."
                    },
                    {
                        "question": "AI가 특정 환경에서 성과에 따른 보상을 받으며 스스로 학습하는 방식은?",
                        "options": ["강화 학습", "지도 학습", "비지도 학습", "고정 학습"],
                        "answer": 1,
                        "explanation": "강화 학습(Reinforcement Learning)은 행동에 따른 보상(Reward)을 최대화하는 방향으로 정책을 학습합니다."
                    }
                ]
            }
            self.after(0, lambda: self.display_results(sample_data))
            self.after(0, lambda: messagebox.showinfo("안내", "API 키가 없어 샘플 시연 모드로 실행되었습니다."))
            return

        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-1.5-flash')
            
            prompt = f"""
            당신은 최고의 교육 도우미 AI입니다. 아래 제공된 강의 텍스트를 분석하여 두 가지를 생성하세요.
            모든 답변은 한국어로 작성하세요.

            1. 강의 핵심 요약: 정확히 5개 문장으로 내용을 요약하세요 (번호 포함).
            2. 예상 문제 생성: 내용을 바탕으로 한 객관식 문제 3개를 만드세요.
               - 각 문제는 4지 선다형이어야 합니다.
               - 문제, 보기, 정답(1~4 중 숫자), 그리고 해설을 반드시 포함하세요.

            결과는 반드시 아래의 JSON 형식을 정확히 지켜서 'JSON_START'와 'JSON_END' 태그 사이에 넣어주세요.
            JSON_START
            {{
              "summary": ["문장1", "문장2", "문장3", "문장4", "문장5"],
              "quizzes": [
                {{
                  "question": "문제 내용",
                  "options": ["보기1", "보기2", "보기3", "보기4"],
                  "answer": 1,
                  "explanation": "상세한 해설 내용"
                }},
                ... (3개)
              ]
            }}
            JSON_END

            강의 내용:
            {content}
            """
            
            response = model.generate_content(prompt)
            result_text = response.text
            
            json_match = re.search(r'JSON_START(.*?)JSON_END', result_text, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group(1).strip())
                self.after(0, lambda: self.display_results(data))
            else:
                raise Exception("AI 응답에서 데이터 형식을 찾을 수 없습니다.")

        except Exception as e:
            self.after(0, lambda: messagebox.showerror("AI 분석 오류", f"오류 발생: {str(e)}"))
        finally:
            self.after(0, lambda: self.process_btn.configure(state="normal", text="✨ AI 요약 및 문제 생성 시작"))

    def display_results(self, data):
        # 1. 요약 표시
        self.summary_textbox.configure(state="normal")
        self.summary_textbox.delete("1.0", tk.END)
        summary_text = "\n\n".join(data['summary'])
        self.summary_textbox.insert("1.0", summary_text)
        self.summary_textbox.configure(state="disabled")

        # 2. 퀴즈 표시
        for widget in self.quiz_scroll.winfo_children():
            widget.destroy()

        self.quiz_data = data['quizzes']
        self.user_answers = []
        self.explanation_labels = []

        for i, q in enumerate(self.quiz_data):
            q_frame = ctk.CTkFrame(self.quiz_scroll, corner_radius=10)
            q_frame.pack(fill="x", padx=10, pady=10)

            label = ctk.CTkLabel(q_frame, text=f"Q{i+1}. {q['question']}", 
                                font=ctk.CTkFont(size=15, weight="bold"),
                                wraplength=700, justify="left")
            label.pack(anchor="w", padx=20, pady=(15, 10))

            v = tk.IntVar(value=0)
            self.user_answers.append(v)

            for idx, opt in enumerate(q['options']):
                rb = ctk.CTkRadioButton(q_frame, text=opt, variable=v, value=idx+1)
                rb.pack(anchor="w", padx=40, pady=5)

            # 해설용 레이블 (초기에는 숨김 또는 빈 상태)
            exp_label = ctk.CTkLabel(q_frame, text="", wraplength=650, justify="left", font=ctk.CTkFont(size=12))
            exp_label.pack(anchor="w", padx=20, pady=(10, 15))
            self.explanation_labels.append(exp_label)

        self.grade_btn.configure(state="normal")
        self.tabview.set("5줄 요약") # 요약 탭으로 자동 이동

    def grade_quiz(self):
        score = 0
        total = len(self.quiz_data)

        for i, (v, q) in enumerate(zip(self.user_answers, self.quiz_data)):
            user_choice = v.get()
            correct = q['answer']
            exp_label = self.explanation_labels[i]

            if user_choice == 0:
                exp_label.configure(text="⚠️ 문제를 풀어주세요!", text_color="#e74c3c")
                continue

            if user_choice == correct:
                score += 1
                exp_label.configure(text=f"✅ 정답입니다!\n{q['explanation']}", text_color="#2ecc71")
            else:
                exp_label.configure(text=f"❌ 틀렸습니다. (정답: {correct}번)\n{q['explanation']}", text_color="#e74c3c")
        
        messagebox.showinfo("결과", f"총 {total}문제 중 {score}문제를 맞혔습니다!")

if __name__ == "__main__":
    app = StudyAssistantApp()
    app.mainloop()
